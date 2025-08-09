# Core logic for handling all the extraction of text from pptx
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

###

def extract_shapes(shapes, collected_data, slide_index):
    for shape in shapes:
        # Group shapes in any pptx
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:  
            extract_shapes(shape.shapes, collected_data, slide_index)
        
        # Extract text content from shapes with a text frame
        if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = ''.join(run.text for run in paragraph.runs)
                if text:
                    #print("Text:", text)
                    collected_data['texts'].append({
                        'slide': slide_index,
                        'text': text
                    })

            # Also extract alternate text if available (+ image alt)
            if hasattr(shape, 'alt_text') and shape.alt_text.strip():
                #print("Alt Text:", shape.alt_text.strip())
                collected_data['alt_texts'].append({
                    'slide': slide_index,
                    'alt_text': shape.alt_text.strip()
                })
        
        # Extract text from tables
        if hasattr(shape, 'has_table') and shape.has_table:
            table_data = []
            table = shape.table
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)

            collected_data['tables'].append({
                'slide': slide_index,
                'table': table_data
            })
            #print("Table text: ", table_data)

        # Extract image_data.
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                image = shape.image
                image_bytes = image.blob
                filename = f"slide{slide_index}_image_{getattr(shape, 'shape_id', 'unknown')}.{image.ext}"  #num-id/u-jpeg type

                image_info = {
                    'slide': slide_index,
                    'filename': filename,
                    'blob': image_bytes,      
                    'ext': image.ext,
                    'content_type': image.content_type,
                    'shape_id': getattr(shape, 'shape_id', None)
                }
                if hasattr(shape, 'alt_text') and shape.alt_text.strip():
                    image_info['alt_text'] = shape.alt_text.strip()

                collected_data['images'].append(image_info)
            except Exception as ex:
                print(f"couldn't parse: {slide_index}: {ex}")

        # Finally chart data(pie, col etc)
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            try:
                chart = shape.chart
                chart_data = {
                    'slide': slide_index,
                    'title': None,
                    'series': []
                }

                try:
                    if chart.chart_title and chart.chart_title.text_frame:
                        chart_data['title'] = chart.chart_title.text_frame.text.strip()
                except Exception:
                    chart_data['title'] = None

                try:
                    for series in chart.series:
                        series_name = getattr(series, 'name', None)
                        values = []
                        for point in series.points:
                            values.append(getattr(point, 'value', None))
                        chart_data['series'].append({'name': series_name, 'values': values})
                except Exception:
                    pass

                chart_data['empty'] = all(
                    all(v is None or v == "" for v in s['values'])
                    for s in chart_data['series']
                )

                collected_data['charts'].append(chart_data)
            except Exception as ex:
                print(f"couldn't read chart on: {slide_index}: {ex}")

###

def parse(file_path):

    collected_data = {
        'file': file_path,
        'slide_titles': [],
        'texts': [],
        'alt_texts': [],
        'tables': [],
        'images': [],  
        'charts': [],
        'notes': [],
    }

    # catch errors incase the path wasn't forwarded.
    try:
        presentation = Presentation(file_path)
        for i, slide in enumerate(presentation.slides, start=1):
            #print(f"\n--- Slide {i+1} ---")
            try:
                # Titles of the files first.
                if slide.shapes.title:
                    title_text = slide.shapes.title.text.strip()
                    if title_text:
                        collected_data['slide_titles'].append({'slide': i, 'title': title_text})
                    else:
                        collected_data['slide_titles'].append({'slide': i, 'title': None})
                else:
                    collected_data['slide_titles'].append({'slide': i, 'title': None})
            except Exception:
                collected_data['slide_titles'].append({'slide': i, 'title': None})
            
            extract_shapes(slide.shapes, collected_data, slide_index=i)
                
            # Extract user notes if they exist.
            if hasattr(slide, 'notes_slide') and slide.notes_slide:
                notes_text_frame = slide.notes_slide.notes_text_frame
                if notes_text_frame:
                    notes_text = notes_text_frame.text.strip()
                    if notes_text:
                        collected_data['notes'].append({'slide': i, 'notes': notes_text})
            
        return collected_data
    
    except Exception as ex:
        print(f"Error in parsing pptx: {ex}")
        return collected_data
